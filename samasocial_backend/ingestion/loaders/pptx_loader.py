# pptx_loader.py
import logging
import tempfile
import os
from pptx import Presentation
from pptx.util import Pt
from langchain_core.documents import Document

log = logging.getLogger(__name__)


def load_pptx(file_bytes: bytes, filename: str) -> list[Document]:
    """
    Extract text from each PPTX slide.
    Preserves slide numbers for citations.
    ASSIGNMENT CITATION: "from slide 4 of filename"
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        docs = []
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []

            for shape in slide.shapes:
                # WHY CHECK HAS_TEXT_FRAME:
                # Not all shapes have text — images, charts don't
                # Checking prevents AttributeError
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = " ".join(
                        run.text for run in para.runs
                        if run.text.strip()
                    )
                    if line.strip():
                        texts.append(line.strip())

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"Notes: {notes}")

            slide_text = "\n".join(texts)

            if slide_text.strip():
                docs.append(Document(
                    page_content=slide_text,
                    metadata={
                        "source_type": "pptx",
                        "source_file": filename,
                        "slide_number": slide_num,
                        "total_slides": total_slides,
                        "citation": f"Slide {slide_num}/{total_slides} of {filename}",
                    }
                ))

        log.info(f"PPTX loaded: {len(docs)} slides from '{filename}'")

        if not docs:
            raise ValueError(f"No text found in PPTX: {filename}")

        return docs

    finally:
        os.unlink(tmp_path)